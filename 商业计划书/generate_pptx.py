"""
生成智径 AdaptivePath 路演 PPT
====================================

精心设计的演示文稿，突出创新性与技术深度。
"""
import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy


# ---------------------------------------------------------------------- #
# 主题色
# ---------------------------------------------------------------------- #
PRIMARY = RGBColor(0x00, 0x75, 0xFF)        # 科技蓝
SECONDARY = RGBColor(0x10, 0xB9, 0x81)      # 活力绿
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)         # 警告橙
DARK = RGBColor(0x1F, 0x29, 0x37)           # 深色
LIGHT = RGBColor(0xF9, 0xFA, 0xFB)          # 浅色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)           # 灰
RED = RGBColor(0xEF, 0x44, 0x44)            # 红
GRADIENT_START = RGBColor(0x00, 0x75, 0xFF)
GRADIENT_END = RGBColor(0x00, 0xC4, 0x9A)


# ---------------------------------------------------------------------- #
# 工具函数
# ---------------------------------------------------------------------- #
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_outline(shape):
    shape.line.fill.background()


def add_text_box(slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    return tb


def add_background(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_fill(bg, color)
    set_no_outline(bg)
    return bg


def add_accent_bar(slide, color=PRIMARY):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.15))
    set_fill(bar, color)
    set_no_outline(bar)
    return bar


def add_section_title(slide, title, subtitle=""):
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), title, size=32, color=DARK, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4), subtitle, size=16, color=GRAY)
    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.04))
    set_fill(line, PRIMARY)
    set_no_outline(line)


def add_card(slide, x, y, w, h, fill=WHITE, shadow=True):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(card, fill)
    set_no_outline(card)
    # 阴影效果（通过 outline 模拟）
    if shadow:
        card.shadow.inherit = False
    return card


def add_icon_circle(slide, x, y, size, color, text, text_color=WHITE, font_size=24):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    set_fill(circle, color)
    set_no_outline(circle)
    tf = circle.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = True
    run.font.name = "Microsoft YaHei"
    return circle


def add_numbered_point(slide, x, y, num, text, w=Inches(11), color=PRIMARY, size=16):
    add_text_box(slide, x, y, Inches(0.6), Inches(0.4), f"{num}.", size=size+2, color=color, bold=True)
    add_text_box(slide, x + Inches(0.6), y, w - Inches(0.6), Inches(0.4), text, size=size, color=DARK)


# ---------------------------------------------------------------------- #
# 1. 创建 PPT
# ---------------------------------------------------------------------- #
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ---------------------------------------------------------------------- #
# Slide 1: 封面
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide, DARK)
# 装饰
for i in range(5):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5 + i * 0.4), Inches(0.3), Inches(0.15), Inches(0.15))
    set_fill(circle, PRIMARY if i % 2 == 0 else SECONDARY)
    set_no_outline(circle)

# Logo 占位
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.5), "🎓 智径 AdaptivePath", size=18, color=PRIMARY, bold=True)

# 标题
add_text_box(slide, Inches(0.5), Inches(2.2), Inches(12), Inches(1.2), "智径 AdaptivePath", size=60, color=WHITE, bold=True)
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(12), Inches(0.8), "基于多智能体协同与认知诊断融合的自适应学习伴学智能体", size=24, color=PRIMARY)
add_text_box(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(0.5), "Multi-Agent Adaptive Learning Companion", size=18, color=GRAY)

# 信息
info = [
    ("赛道", "科大讯飞 AI 开发者大赛 · 自适应学习路径决策与伴学智能体"),
    ("作品类型", "AI Agent + 教育认知科学"),
    ("团队", "崇理团队"),
    ("日期", "2026 年 8 月"),
]
for i, (k, v) in enumerate(info):
    y = Inches(5.5 + i * 0.4)
    add_text_box(slide, Inches(0.5), y, Inches(2), Inches(0.3), k + "：", size=14, color=GRAY, bold=True)
    add_text_box(slide, Inches(2.5), y, Inches(10), Inches(0.3), v, size=14, color=WHITE)


# ---------------------------------------------------------------------- #
# Slide 2: 目录
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)

add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6), "目录", size=32, color=DARK, bold=True)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4), "CONTENTS", size=14, color=GRAY)

contents = [
    ("01", "项目背景与洞察", "为什么需要重新思考 AI 教育"),
    ("02", "核心创新点", "CSN 三模型融合 + Bandit + Socratic"),
    ("03", "系统架构", "5 智能体协同架构"),
    ("04", "技术实现", "代码 + 公式 + 流程"),
    ("05", "产品演示", "Web Demo 实操展示"),
    ("06", "市场与商业", "TAM/SAM/SOM + 商业模式"),
    ("07", "团队与未来", "团队 + 路线图"),
]

for i, (num, title, desc) in enumerate(contents):
    y = Inches(1.8 + i * 0.7)
    add_text_box(slide, Inches(0.8), y, Inches(0.8), Inches(0.5), num, size=28, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(4), Inches(0.4), title, size=20, color=DARK, bold=True)
    add_text_box(slide, Inches(6), y + Inches(0.05), Inches(7), Inches(0.4), desc, size=14, color=GRAY)


# ---------------------------------------------------------------------- #
# Slide 3: 项目背景
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "01 项目背景", "教育大模型浪潮中的真实痛点")

# 大数据
stats = [
    ("5000亿", "中国 K12 在线教育市场规模（元）"),
    ("90%", "现有 AI 教育产品仅做 LLM 套壳"),
    ("0", "真正拥有学生认知模型的 AI 教育产品"),
    ("200+", "教育认知科学可借鉴的方法"),
]
for i, (num, desc) in enumerate(stats):
    x = Inches(0.5 + i * 3.1)
    add_card(slide, x, Inches(1.8), Inches(2.9), Inches(1.8), fill=LIGHT)
    add_text_box(slide, x + Inches(0.2), Inches(1.95), Inches(2.5), Inches(0.7), num, size=42, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.8), desc, size=12, color=DARK, align=PP_ALIGN.CENTER)

# 四大痛点
add_text_box(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.5), "现有 AI 教育产品的四大通病", size=22, color=DARK, bold=True)

pains = [
    ("❌", "无学情诊断", "把学生当答题机器，不了解掌握度"),
    ("❌", "路径千篇一律", "知识图谱推荐，不考虑学生差异"),
    ("❌", "灌输式教学", "直接给答案，违背 ZPD 原则"),
    ("❌", "忽视情感", "挫败时继续推难题，体验差"),
]
for i, (icon, title, desc) in enumerate(pains):
    x = Inches(0.5 + i * 3.1)
    add_card(slide, x, Inches(4.7), Inches(2.9), Inches(2.2), fill=WHITE)
    add_text_box(slide, x + Inches(0.2), Inches(4.85), Inches(0.6), Inches(0.5), icon, size=28, color=RED)
    add_text_box(slide, x + Inches(0.2), Inches(5.4), Inches(2.5), Inches(0.4), title, size=18, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(5.85), Inches(2.5), Inches(0.9), desc, size=12, color=GRAY)


# ---------------------------------------------------------------------- #
# Slide 4: 核心创新点 - 大图
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "02 核心创新点", "五大维度全面超越 LLM 套壳")

# 5 个创新卡片
innovations = [
    ("🎯", "CSN", "三模型融合", "DKT × BKT × IRT", "样本量自适应", PRIMARY),
    ("🎰", "Bandit", "在线路径优化", "Contextual Bandit", "冷启动友好", SECONDARY),
    ("🎓", "Socratic", "苏格拉底式教学", "5 级脚手架", "Bloom 适配", ACCENT),
    ("❤️", "Affective", "情感计算", "5 维实时感知", "影响教学节奏", RED),
    ("🔍", "XAI", "完整可解释性", "每次推荐有理由", "可追溯", PRIMARY),
]
for i, (icon, name, en, desc1, desc2, color) in enumerate(innovations):
    x = Inches(0.4 + i * 2.55)
    add_card(slide, x, Inches(2.0), Inches(2.4), Inches(4.8), fill=LIGHT)
    # 顶部色条
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.4), Inches(0.1))
    set_fill(top, color)
    set_no_outline(top)
    # Icon
    add_text_box(slide, x + Inches(0.8), Inches(2.3), Inches(0.8), Inches(0.8), icon, size=42, color=color, align=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide, x + Inches(0.1), Inches(3.3), Inches(2.2), Inches(0.5), name, size=20, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(3.8), Inches(2.2), Inches(0.4), en, size=12, color=color, align=PP_ALIGN.CENTER, bold=True)
    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.6), Inches(4.3), Inches(1.2), Inches(0.03))
    set_fill(div, color)
    set_no_outline(div)
    # Descs
    add_text_box(slide, x + Inches(0.2), Inches(4.5), Inches(2.0), Inches(0.4), desc1, size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(5.0), Inches(2.0), Inches(0.4), desc2, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # 关键数字
    add_text_box(slide, x + Inches(0.2), Inches(5.8), Inches(2.0), Inches(0.6),
                ["5 模型", "32 维", "5 阶段", "5 维度", "100%"][i],
                size=24, color=color, bold=True, align=PP_ALIGN.CENTER)

# 底部核心声明
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
             "💎 不是堆叠名词，是工程级深度融合：每行代码都对应一篇经典论文", size=14, color=DARK, align=PP_ALIGN.CENTER, bold=True)


# ---------------------------------------------------------------------- #
# Slide 5: CSN 三模型融合详解
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "创新 1：CSN 三模型融合", "BKT × DKT × IRT 自适应加权融合")

# 左侧：三个模型卡片
models = [
    ("BKT", "Bayesian Knowledge Tracing", "4 状态 HMM，冷启动友好", "P(L_t | obs) = ?", PRIMARY),
    ("DKT", "Deep Knowledge Tracing", "LSTM，长程记忆", "h_t = LSTM(x_t, h_{t-1})", SECONDARY),
    ("IRT", "Item Response Theory 2PL", "全局能力估计", "P(θ) = σ(a(θ-b))", ACCENT),
]
for i, (name, en, desc, formula, color) in enumerate(models):
    y = Inches(2.0 + i * 1.4)
    add_card(slide, Inches(0.5), y, Inches(5.5), Inches(1.2), fill=LIGHT)
    add_icon_circle(slide, Inches(0.7), y + Inches(0.3), Inches(0.6), color, name[:1], text_color=WHITE, font_size=20)
    add_text_box(slide, Inches(1.5), y + Inches(0.15), Inches(4), Inches(0.4), name, size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.5), Inches(4), Inches(0.3), en, size=10, color=GRAY)
    add_text_box(slide, Inches(1.5), y + Inches(0.75), Inches(4), Inches(0.4), desc, size=12, color=DARK)
    add_text_box(slide, Inches(5.0), y + Inches(0.3), Inches(2.5), Inches(0.6), formula, size=12, color=color, bold=True, align=PP_ALIGN.RIGHT)

# 右侧：融合示意
add_text_box(slide, Inches(6.5), Inches(2.0), Inches(6.5), Inches(0.5), "自适应加权融合", size=20, color=DARK, bold=True)

# 融合公式
add_card(slide, Inches(6.5), Inches(2.7), Inches(6.5), Inches(1.3), fill=LIGHT)
add_text_box(slide, Inches(6.7), Inches(2.85), Inches(6.1), Inches(0.4), "mastery = w_dkt·dkt_p + w_bkt·bkt_p + w_irt·irt_p", size=16, color=DARK, bold=True, font="Consolas")
add_text_box(slide, Inches(6.7), Inches(3.25), Inches(6.1), Inches(0.4), "confidence = 1 / (1 + 10·Var([dkt_p, bkt_p, irt_p]))", size=14, color=DARK, font="Consolas")
add_text_box(slide, Inches(6.7), Inches(3.65), Inches(6.1), Inches(0.3), "(置信度基于三模型方差，样本量自适应)", size=11, color=GRAY)

# 权重自适应
add_text_box(slide, Inches(6.5), Inches(4.2), Inches(6.5), Inches(0.4), "样本量自适应权重", size=16, color=DARK, bold=True)

weights = [
    ("n<5", "(0.2, 0.5, 0.3)", "BKT 主导\n冷启动", PRIMARY),
    ("5≤n<20", "(0.4, 0.3, 0.3)", "三模型\n等权", SECONDARY),
    ("n≥20", "(0.5, 0.2, 0.3)", "DKT 主导\n长程依赖", ACCENT),
]
for i, (cond, w, label, color) in enumerate(weights):
    x = Inches(6.5 + i * 2.2)
    add_card(slide, x, Inches(4.7), Inches(2.0), Inches(2.0), fill=WHITE)
    add_text_box(slide, x, Inches(4.85), Inches(2.0), Inches(0.4), cond, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.3), Inches(2.0), Inches(0.3), w, size=11, color=DARK, align=PP_ALIGN.CENTER, font="Consolas")
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.4), Inches(5.7), Inches(1.2), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x, Inches(5.85), Inches(2.0), Inches(0.8), label, size=12, color=DARK, align=PP_ALIGN.CENTER, bold=True)

# 底部
add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5),
             "✨ 业内首次：将三种认知诊断模型在同一系统内自适应融合 + 置信度估计", size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 6: Contextual Bandit 路径优化
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "创新 2：Contextual Bandit 路径优化", "在线学习最优推荐策略")

# 左侧：问题与方案
add_text_box(slide, Inches(0.5), Inches(2.0), Inches(6), Inches(0.5), "🎯 核心问题", size=20, color=DARK, bold=True)
add_text_box(slide, Inches(0.5), Inches(2.5), Inches(6), Inches(0.5), "如何为每个学生在线学习最优的下一步推荐？", size=14, color=DARK)

add_text_box(slide, Inches(0.5), Inches(3.2), Inches(6), Inches(0.5), "💡 方案：Contextual Bandit", size=20, color=PRIMARY, bold=True)

points = [
    ("Context (上下文)", "x ∈ R³²: 学生状态向量"),
    ("Action (动作)", "a ∈ A: 候选知识点"),
    ("Reward (奖励)", "r: 学习收益"),
    ("Algorithm", "LinUCB (Linear Upper Confidence Bound)"),
]
for i, (k, v) in enumerate(points):
    y = Inches(3.8 + i * 0.4)
    add_text_box(slide, Inches(0.7), y, Inches(2), Inches(0.3), k, size=13, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(2.7), y, Inches(4), Inches(0.3), v, size=13, color=DARK)

# LinUCB 公式
add_card(slide, Inches(0.5), Inches(5.7), Inches(6.0), Inches(1.5), fill=LIGHT)
add_text_box(slide, Inches(0.7), Inches(5.85), Inches(5.6), Inches(0.4), "LinUCB 选择规则", size=14, color=PRIMARY, bold=True)
add_text_box(slide, Inches(0.7), Inches(6.25), Inches(5.6), Inches(0.5), "a* = argmax_a [ θ_a^T x + α √(x^T A_a⁻¹ x) ]", size=16, color=DARK, bold=True, font="Consolas")
add_text_box(slide, Inches(0.7), Inches(6.75), Inches(5.6), Inches(0.4), "利用 + 探索 + 理论保证（O(√T) regret bound）", size=11, color=GRAY)

# 右侧：ZPD 融合
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(6), Inches(0.5), "🌱 ZPD 融合（Vygotsky）", size=20, color=DARK, bold=True)

# ZPD 公式
add_card(slide, Inches(7.0), Inches(2.7), Inches(6.0), Inches(1.5), fill=LIGHT)
add_text_box(slide, Inches(7.2), Inches(2.85), Inches(5.6), Inches(0.4), "ZPD 分数", size=14, color=SECONDARY, bold=True)
add_text_box(slide, Inches(7.2), Inches(3.25), Inches(5.6), Inches(0.5), "zpd(θ, d) = exp(-((d - θ - 0.4)² / 0.2))", size=14, color=DARK, font="Consolas")
add_text_box(slide, Inches(7.2), Inches(3.75), Inches(5.6), Inches(0.4), "差距 0.4：略感挑战但可完成", size=11, color=GRAY)

# 优势
add_text_box(slide, Inches(7.0), Inches(4.5), Inches(6), Inches(0.5), "💪 优势", size=20, color=DARK, bold=True)

advantages = [
    "✅ 冷启动友好 - 无需离线训练数据",
    "✅ 在线学习 - 每步可更新",
    "✅ 理论保证 - regret bound",
    "✅ ZPD 过滤 - 避免过易/过难",
    "✅ 可解释 - 每个推荐都有 UCB 分数",
]
for i, a in enumerate(advantages):
    y = Inches(5.1 + i * 0.35)
    add_text_box(slide, Inches(7.2), y, Inches(6), Inches(0.3), a, size=12, color=DARK)


# ---------------------------------------------------------------------- #
# Slide 7: 苏格拉底式教学
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "创新 3：苏格拉底式支架教学", "永远不直接给答案，但会渐进逼近")

# 5 阶段状态机
add_text_box(slide, Inches(0.5), Inches(1.9), Inches(12), Inches(0.4), "5 阶段对话状态机", size=18, color=DARK, bold=True)

stages = [
    ("DIAGNOSE", "了解认知"),
    ("PROBE", "反诘深入"),
    ("HINT", "关键提示"),
    ("CONFIRM", "检验理解"),
    ("REFLECT", "元反思"),
]
for i, (s, d) in enumerate(stages):
    x = Inches(0.5 + i * 2.55)
    add_card(slide, x, Inches(2.4), Inches(2.4), Inches(1.2), fill=LIGHT)
    add_text_box(slide, x, Inches(2.5), Inches(2.4), Inches(0.5), s, size=16, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.05), Inches(2.4), Inches(0.5), d, size=13, color=DARK, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.85 + i * 2.55), Inches(2.85), Inches(0.2), Inches(0.3))
        set_fill(arrow, PRIMARY)
        set_no_outline(arrow)

# 5 级脚手架
add_text_box(slide, Inches(0.5), Inches(3.9), Inches(12), Inches(0.4), "5 级动态脚手架（基于 Bloom 认知分类）", size=18, color=DARK, bold=True)

scaffoldings = [
    ("L0", "元认知提示", "掌握度≥0.85", "你想想看，第一步应该是什么？"),
    ("L1", "关键概念", "0.7≤掌握度<0.85", "提示：这与【梯度下降】相关"),
    ("L2", "类比示例", "0.5≤掌握度<0.7", "想象你在下山..."),
    ("L3", "分步分解", "0.3≤掌握度<0.5", "第 1 步...第 2 步..."),
    ("L4", "完整示范", "掌握度<0.3", "演示完整解法"),
]
for i, (lvl, name, cond, example) in enumerate(scaffoldings):
    x = Inches(0.5 + i * 2.55)
    color = [PRIMARY, SECONDARY, ACCENT, RGBColor(0xEF, 0x44, 0x44), RGBColor(0x8B, 0x5C, 0xF6)][i]
    add_card(slide, x, Inches(4.4), Inches(2.4), Inches(2.6), fill=WHITE)
    # level indicator
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), Inches(4.55), Inches(0.6), Inches(0.6))
    set_fill(badge, color)
    set_no_outline(badge)
    tf = badge.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = lvl
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Microsoft YaHei"
    add_text_box(slide, x + Inches(0.9), Inches(4.55), Inches(1.4), Inches(0.4), name, size=14, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.9), Inches(4.85), Inches(1.4), Inches(0.3), cond, size=9, color=GRAY)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), Inches(5.3), Inches(2.0), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x + Inches(0.2), Inches(5.4), Inches(2.0), Inches(1.5), f"「{example}」", size=11, color=DARK, align=PP_ALIGN.CENTER)

# 核心原则
add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.4),
             "🎓 核心原则：连续 2 轮错→升级 | 连续 2 轮对→降级 | 主动求提示→中等 | 疲劳→降级", size=12, color=DARK, align=PP_ALIGN.CENTER, bold=True)


# ---------------------------------------------------------------------- #
# Slide 8: 系统架构图
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "03 系统架构", "5 智能体协同 + 底层支撑")

# 顶部：主控器
add_card(slide, Inches(3.5), Inches(2.0), Inches(6.3), Inches(0.8), fill=PRIMARY)
add_text_box(slide, Inches(3.5), Inches(2.1), Inches(6.3), Inches(0.6), "🎯 Orchestrator Agent (主控调度)", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 5 个智能体
agents = [
    ("🎯", "Diagnostic", PRIMARY),
    ("🗺️", "Planning", SECONDARY),
    ("🎓", "Teaching", ACCENT),
    ("💭", "Reflective", RGBColor(0x8B, 0x5C, 0xF6)),
    ("❤️", "Emotional", RED),
]
for i, (icon, name, color) in enumerate(agents):
    x = Inches(0.4 + i * 2.6)
    add_card(slide, x, Inches(3.3), Inches(2.4), Inches(1.0), fill=LIGHT)
    add_text_box(slide, x, Inches(3.4), Inches(2.4), Inches(0.4), icon, size=20, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.8), Inches(2.4), Inches(0.4), name, size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)

# 底层：CSN + Memory
add_text_box(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4), "底层支撑", size=14, color=GRAY, align=PP_ALIGN.CENTER, bold=True)

add_card(slide, Inches(0.5), Inches(5.1), Inches(3.9), Inches(0.8), fill=WHITE)
add_text_box(slide, Inches(0.5), Inches(5.2), Inches(3.9), Inches(0.4), "🧠 Cognitive State Network", size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(5.55), Inches(3.9), Inches(0.3), "DKT × BKT × IRT 融合", size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_card(slide, Inches(4.7), Inches(5.1), Inches(3.9), Inches(0.8), fill=WHITE)
add_text_box(slide, Inches(4.7), Inches(5.2), Inches(3.9), Inches(0.4), "💾 Long-term Memory", size=14, color=SECONDARY, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.7), Inches(5.55), Inches(3.9), Inches(0.3), "Episodic + Semantic + Reflective", size=11, color=GRAY, align=PP_ALIGN.CENTER)

add_card(slide, Inches(8.9), Inches(5.1), Inches(3.9), Inches(0.8), fill=WHITE)
add_text_box(slide, Inches(8.9), Inches(5.2), Inches(3.9), Inches(0.4), "🕸️ Knowledge Graph + RAG", size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.9), Inches(5.55), Inches(3.9), Inches(0.3), "21 节点 + 向量检索", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# 闭环说明
add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4), "完整闭环", size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), "学生输入 → 情感检测 → 学情更新 → 路径规划 → 教学回应 → 反思记录 → (回到学生)", size=13, color=PRIMARY, align=PP_ALIGN.CENTER, bold=True)


# ---------------------------------------------------------------------- #
# Slide 9: 知识图谱
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "知识图谱", "面向 AI 专业导论的精细化图谱")

categories = [
    ("数学基础", 4, ["线性代数", "微积分", "概率论", "最优化"], PRIMARY),
    ("机器学习", 5, ["基本概念", "线性回归", "逻辑回归", "决策树", "SVM/集成"], SECONDARY),
    ("深度学习", 5, ["感知机", "反向传播", "CNN", "RNN", "Transformer"], ACCENT),
    ("应用", 3, ["NLP基础", "预训练模型", "计算机视觉"], RGBColor(0x8B, 0x5C, 0xF6)),
    ("前沿", 2, ["强化学习", "LLM Agent"], RED),
    ("伦理", 1, ["AI伦理"], RGBColor(0xEC, 0x48, 0x99)),
]
for i, (cat, n, skills, color) in enumerate(categories):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.0 + row * 2.4)
    add_card(slide, x, y, Inches(4.0), Inches(2.2), fill=LIGHT)
    # 顶部色条
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(4.0), Inches(0.1))
    set_fill(top, color)
    set_no_outline(top)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.4), f"{cat} ({n})", size=18, color=color, bold=True)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.7), Inches(3.6), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    skills_text = " · ".join(skills)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.6), Inches(1.2), skills_text, size=13, color=DARK)

# 关系类型
add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4), "关系类型：prerequisite（先决）+ related（相关）+ confused_with（易混淆）", size=13, color=DARK, align=PP_ALIGN.CENTER, bold=True)
add_text_box(slide, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.3), "每节点包含：难度、估计学习时长、Bloom 层级、关键概念、典型错误", size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 10: 产品演示（Web Demo）
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "05 产品演示", "Streamlit Web Demo 实时交互")

# 描述
add_text_box(slide, Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.4),
             "完整 Web Demo 包含 6 个页面：对话 / 学情诊断 / 路径规划 / 学习报告 / 智能体架构 / 关于", size=16, color=DARK, align=PP_ALIGN.CENTER)

# 4 个核心功能截图占位
features = [
    ("🏠", "主页对话", "苏格拉底式对话\n5 级脚手架", PRIMARY),
    ("📊", "学情诊断", "CSN 三模型\n置信度可视化", SECONDARY),
    ("🗺️", "路径规划", "Bandit 推荐\n知识图谱可视化", ACCENT),
    ("📝", "学习报告", "情感 + 反思\n元认知培养", RGBColor(0x8B, 0x5C, 0xF6)),
]
for i, (icon, name, desc, color) in enumerate(features):
    x = Inches(0.5 + i * 3.15)
    add_card(slide, x, Inches(2.6), Inches(2.95), Inches(3.5), fill=WHITE)
    # 顶部色条
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.6), Inches(2.95), Inches(0.1))
    set_fill(top, color)
    set_no_outline(top)
    add_text_box(slide, x, Inches(2.9), Inches(2.95), Inches(0.8), icon, size=42, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.8), Inches(2.95), Inches(0.4), name, size=18, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.6), Inches(4.3), Inches(1.75), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x, Inches(4.4), Inches(2.95), Inches(1.5), desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
             "运行命令：streamlit run demo/app.py", size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             "✓ 5 智能体协同 ✓ 21 节点知识图谱 ✓ 25+ 题库 ✓ 端到端闭环", size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 11: 演示流程（实操）
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "演示流程", "从启动会话到反思总结的完整闭环")

steps = [
    ("Step 1", "启动会话", "选择学习目标：llm_agent", "系统生成 5 步学习路径", PRIMARY),
    ("Step 2", "学生提问", "“我不懂 ReAct”", "情感检测：engagement 上升", SECONDARY),
    ("Step 3", "智能回应", "苏格拉底式反诘", "脚手架 L4 → 完整示范", ACCENT),
    ("Step 4", "答题诊断", "答对率 80%", "CSN 提升掌握度到 0.7", RGBColor(0x8B, 0x5C, 0xF6)),
    ("Step 5", "动态规划", "Bandit 推荐下一步", "ZPD 区间内 + 路径优化", RED),
    ("Step 6", "反思总结", "生成学习报告", "强项 / 弱项 / 下一步", PRIMARY),
]

for i, (step, title, input_, output, color) in enumerate(steps):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.0 + row * 2.4)
    add_card(slide, x, y, Inches(4.0), Inches(2.2), fill=LIGHT)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2), Inches(0.4), step, size=14, color=color, bold=True)
    add_text_box(slide, x + Inches(1.0), y + Inches(0.15), Inches(2.8), Inches(0.4), title, size=16, color=DARK, bold=True)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.7), Inches(3.6), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(3.6), Inches(0.5), f"输入：{input_}", size=12, color=DARK)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.6), Inches(0.6), f"输出：{output}", size=12, color=color, bold=True)


# ---------------------------------------------------------------------- #
# Slide 12: 评测对比
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "评测对比", "智径 vs 传统 LLM 套壳")

# 表头
headers = ["维度", "传统 LLM 套壳", "智径 AdaptivePath"]
header_y = Inches(2.0)
header_w = [Inches(2.5), Inches(4.5), Inches(5.5)]
xs = [Inches(0.5), Inches(3.0), Inches(7.5)]
for i, h in enumerate(headers):
    add_card(slide, xs[i], header_y, header_w[i], Inches(0.6), fill=PRIMARY)
    add_text_box(slide, xs[i], header_y, header_w[i], Inches(0.6), h, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 数据
data = [
    ("学情诊断", "❌ 无", "✅ CSN 三模型融合"),
    ("路径规划", "❌ 静态或无", "✅ Bandit 在线学习"),
    ("教学方式", "❌ 灌输式", "✅ 苏格拉底式反诘"),
    ("情感适配", "❌ 忽视", "✅ 5 维实时感知"),
    ("可解释性", "❌ 黑盒", "✅ 每次推荐有理由"),
    ("长期记忆", "❌ 短期对话", "✅ 3 层记忆"),
    ("评测指标", "❌ 无", "✅ 多维报告"),
    ("教学理论", "❌ 无", "✅ ZPD + Bloom + Scaffolding"),
]
for i, (dim, bad, good) in enumerate(data):
    y = Inches(2.7 + i * 0.5)
    add_card(slide, Inches(0.5), y, Inches(2.5), Inches(0.45), fill=LIGHT)
    add_card(slide, Inches(3.0), y, Inches(4.5), Inches(0.45), fill=WHITE)
    add_card(slide, Inches(7.5), y, Inches(5.5), Inches(0.45), fill=WHITE)
    add_text_box(slide, Inches(0.5), y, Inches(2.5), Inches(0.45), dim, size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.0), y, Inches(4.5), Inches(0.45), bad, size=12, color=RED, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.5), y, Inches(5.5), Inches(0.45), good, size=12, color=SECONDARY, bold=True, align=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
             "💎 智径 AdaptivePath 在 8 个维度全面超越传统方案", size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 13: 商业化
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "06 市场与商业", "TAM/SAM/SOM + 多元商业模式")

# TAM/SAM/SOM
add_text_box(slide, Inches(0.5), Inches(1.9), Inches(6), Inches(0.5), "🎯 市场空间", size=20, color=DARK, bold=True)

market = [
    ("TAM", "360 亿元", "中国 K12 + 高校 STEM 学生 × 300 元/年", PRIMARY),
    ("SAM", "108 亿元", "STEM 30% 渗透率", SECONDARY),
    ("SOM", "5400 万元", "前 3 年 0.5% 份额", ACCENT),
]
for i, (k, v, d, color) in enumerate(market):
    y = Inches(2.5 + i * 1.3)
    add_card(slide, Inches(0.5), y, Inches(5.5), Inches(1.1), fill=LIGHT)
    add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(1.5), Inches(0.8), k, size=32, color=color, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.1), Inches(3.6), Inches(0.5), v, size=24, color=DARK, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.6), Inches(3.6), Inches(0.4), d, size=11, color=GRAY)

# 商业模式
add_text_box(slide, Inches(6.5), Inches(1.9), Inches(6), Inches(0.5), "💰 商业模式", size=20, color=DARK, bold=True)

models = [
    ("B2B 机构", "60%", "5-20万/校/年", PRIMARY),
    ("B2C 订阅", "30%", "29元/月, 299元/年", SECONDARY),
    ("B2B2C 预装", "10%", "硬件厂商分成", ACCENT),
]
for i, (m, pct, price, color) in enumerate(models):
    y = Inches(2.5 + i * 1.3)
    add_card(slide, Inches(6.5), y, Inches(6.3), Inches(1.1), fill=LIGHT)
    add_text_box(slide, Inches(6.7), y + Inches(0.1), Inches(2.5), Inches(0.8), m, size=20, color=color, bold=True)
    add_text_box(slide, Inches(9.3), y + Inches(0.1), Inches(3.4), Inches(0.5), pct, size=28, color=color, bold=True)
    add_text_box(slide, Inches(9.3), y + Inches(0.6), Inches(3.4), Inches(0.4), price, size=11, color=GRAY)

# 财务预测
add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
             "📈 财务预测：2026 收入 50万 / 2027 收入 500万 / 2028 收入 3000万", size=14, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 14: 团队
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide)
add_accent_bar(slide)
add_section_title(slide, "07 团队与未来", "高校背景 + 多学科交叉")

# 团队成员
members = [
    ("👨‍💼", "主负责人", "整体战略与商务", "高校 / 商科", PRIMARY),
    ("👨‍💻", "技术负责人", "核心架构", "高校 / CS", SECONDARY),
    ("🧠", "AI 算法负责人", "认知诊断 + RL", "高校 / AI", ACCENT),
    ("📚", "教育学顾问", "认知科学理论", "高校 / 教育", RGBColor(0x8B, 0x5C, 0xF6)),
    ("🎨", "产品负责人", "用户体验", "高校 / 设计", RED),
]
for i, (avatar, role, area, bg, color) in enumerate(members):
    x = Inches(0.4 + i * 2.55)
    add_card(slide, x, Inches(2.0), Inches(2.4), Inches(2.5), fill=LIGHT)
    add_text_box(slide, x, Inches(2.2), Inches(2.4), Inches(0.8), avatar, size=42, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.0), Inches(2.4), Inches(0.4), role, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.4), Inches(3.5), Inches(1.6), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x + Inches(0.1), Inches(3.6), Inches(2.2), Inches(0.4), area, size=11, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(4.0), Inches(2.2), Inches(0.4), bg, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# 路线图
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4), "🗺️ 发展路线图", size=18, color=DARK, bold=True, align=PP_ALIGN.CENTER)

phases = [
    ("短期 3M", "扩展学科：数学/物理", "与 1-2 所中学合作试用", PRIMARY),
    ("中期 6M", "B2B 商务拓展", "多模态支持（图像/语音）", SECONDARY),
    ("长期 12M", "海外拓展", "教师端 Dashboard", ACCENT),
]
for i, (phase, g1, g2, color) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    add_card(slide, x, Inches(5.4), Inches(4.0), Inches(1.6), fill=WHITE)
    add_text_box(slide, x, Inches(5.5), Inches(4.0), Inches(0.4), phase, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.2), Inches(5.95), Inches(1.6), Inches(0.02))
    set_fill(div, color)
    set_no_outline(div)
    add_text_box(slide, x + Inches(0.2), Inches(6.1), Inches(3.6), Inches(0.4), "✓ " + g1, size=12, color=DARK, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(6.5), Inches(3.6), Inches(0.4), "✓ " + g2, size=12, color=DARK, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# Slide 15: 总结
# ---------------------------------------------------------------------- #
slide = prs.slides.add_slide(blank)
add_background(slide, DARK)

# 大标题
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.0), "THANK YOU", size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.5), "感谢聆听 · 期待与您深入交流", size=20, color=PRIMARY, align=PP_ALIGN.CENTER)

# 关键数字
stats = [
    ("5", "智能体"),
    ("3", "认知模型"),
    ("21", "知识节点"),
    ("25+", "题库"),
    ("100+", "代码文件"),
    ("2000+", "代码行数"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.5 + i * 2.1)
    add_text_box(slide, x, Inches(3.2), Inches(2.0), Inches(0.8), num, size=42, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(4.0), Inches(2.0), Inches(0.4), label, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# 一句话总结
add_text_box(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.6),
             "智径 AdaptivePath：让每个学生都拥有一个『既懂知识又懂学生』的 AI 导师", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 联系信息
add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             "GitHub: github.com/logicore-code/zhijing-adaptive-path    ·    GitCode: gitcode.com/jiangzeyu-2026/zhijing-adaptive-path    ·    团队: 崇理团队", size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------- #
# 保存
# ---------------------------------------------------------------------- #
output_path = os.path.join(os.path.dirname(__file__), "智径-AdaptivePath-路演PPT.pptx")
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
